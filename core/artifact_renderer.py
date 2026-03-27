"""
core/artifact_renderer.py — Artifact JSON Spec → 실제 파일 변환

PPT → .pptx (python-pptx)
PDF → .pdf (markdown → HTML → weasyprint, fallback: reportlab)
Doc → .md (마크다운 직접 출력)

사용법:
    from core.artifact_renderer import render_to_file
    path = render_to_file(spec, artifact_type="ppt", output_dir="./output")
"""

from __future__ import annotations

import json
import os
from datetime import datetime
from pathlib import Path
from typing import Optional


def render_to_file(
    spec: dict,
    artifact_type: str,
    output_dir: str = "",
    filename: str = "",
) -> str:
    """artifact spec → 실제 파일 생성. 파일 경로 반환."""
    if not output_dir:
        output_dir = str(Path(__file__).parent.parent / "output")
    Path(output_dir).mkdir(parents=True, exist_ok=True)

    if not filename:
        title = spec.get("deck_title", spec.get("document_title", "horcrux_output"))
        safe_title = "".join(c if c.isalnum() or c in "-_ " else "" for c in title)[:50].strip()
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{safe_title}_{ts}"

    if artifact_type == "ppt":
        return _render_pptx(spec, output_dir, filename)
    elif artifact_type == "pdf":
        return _render_pdf(spec, output_dir, filename)
    elif artifact_type == "doc":
        return _render_markdown(spec, output_dir, filename)
    elif artifact_type == "readme":
        return _render_markdown(spec, output_dir, filename, ext="md")
    else:
        return _render_markdown(spec, output_dir, filename)


# ═══════════════════════════════════════════
#  PPT → .pptx
# ═══════════════════════════════════════════

def _render_pptx(spec: dict, output_dir: str, filename: str) -> str:
    """JSON spec → .pptx 파일."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        # python-pptx 없으면 마크다운 fallback
        return _render_ppt_as_markdown(spec, output_dir, filename)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    deck_title = spec.get("deck_title", "Presentation")
    slides_data = spec.get("slides", [])

    for slide_data in slides_data:
        slide_type = slide_data.get("type", "content")
        title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        content = slide_data.get("content", "")
        speaker_note = slide_data.get("speaker_note", "")

        # 레이아웃 선택
        if slide_type == "title":
            layout = prs.slide_layouts[0]  # Title Slide
        elif slide_type == "closing":
            layout = prs.slide_layouts[0]
        else:
            layout = prs.slide_layouts[1]  # Title and Content

        slide = prs.slides.add_slide(layout)

        # 제목
        if slide.shapes.title:
            slide.shapes.title.text = title

        # 내용
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            if bullets:
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        tf.text = str(bullet)
                    else:
                        p = tf.add_paragraph()
                        p.text = str(bullet)
                        p.level = 0
            elif content:
                tf.text = str(content)

        # 스피커 노트
        if speaker_note:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_note

    path = os.path.join(output_dir, f"{filename}.pptx")
    prs.save(path)
    return path


def _render_ppt_as_markdown(spec: dict, output_dir: str, filename: str) -> str:
    """python-pptx 없을 때 마크다운으로 PPT 스펙 출력."""
    lines = [f"# {spec.get('deck_title', 'Presentation')}\n"]
    for slide in spec.get("slides", []):
        lines.append(f"\n---\n## Slide {slide.get('slide_no', '?')}: {slide.get('title', '')}\n")
        for b in slide.get("bullets", []):
            lines.append(f"- {b}")
        if slide.get("content"):
            lines.append(f"\n{slide['content']}")
        if slide.get("speaker_note"):
            lines.append(f"\n> Speaker Note: {slide['speaker_note']}")

    path = os.path.join(output_dir, f"{filename}_slides.md")
    with open(path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    return path


# ═══════════════════════════════════════════
#  PDF → .pdf
# ═══════════════════════════════════════════

def _render_pdf(spec: dict, output_dir: str, filename: str) -> str:
    """JSON spec → .pdf 파일. markdown → HTML → weasyprint, fallback to reportlab."""
    md_content = _spec_to_markdown(spec)

    # 방법 1: weasyprint
    try:
        from weasyprint import HTML
        html = _markdown_to_html(md_content, spec.get("document_title", "Document"))
        path = os.path.join(output_dir, f"{filename}.pdf")
        HTML(string=html).write_pdf(path)
        return path
    except ImportError:
        pass

    # 방법 2: reportlab
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import inch
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        path = os.path.join(output_dir, f"{filename}.pdf")
        doc = SimpleDocTemplate(path, pagesize=A4)
        styles = getSampleStyleSheet()

        # 한글 폰트 시도
        try:
            font_path = _find_korean_font()
            if font_path:
                pdfmetrics.registerFont(TTFont("Korean", font_path))
                styles["Normal"].fontName = "Korean"
                styles["Heading1"].fontName = "Korean"
                styles["Heading2"].fontName = "Korean"
        except Exception:
            pass

        story = []
        title = spec.get("document_title", "Document")
        story.append(Paragraph(title, styles["Title"]))
        story.append(Spacer(1, 0.3 * inch))

        for section in spec.get("sections", []):
            heading = section.get("heading", "")
            content = section.get("content", "")
            if heading:
                story.append(Paragraph(heading, styles["Heading2"]))
                story.append(Spacer(1, 0.1 * inch))
            if content:
                for para in content.split("\n\n"):
                    para = para.strip()
                    if para:
                        # bullet 처리
                        if para.startswith("- ") or para.startswith("* "):
                            for line in para.split("\n"):
                                line = line.strip().lstrip("-* ").strip()
                                if line:
                                    story.append(Paragraph(f"• {line}", styles["Normal"]))
                        else:
                            story.append(Paragraph(para.replace("\n", "<br/>"), styles["Normal"]))
                        story.append(Spacer(1, 0.05 * inch))

        doc.build(story)
        return path
    except ImportError:
        pass

    # 방법 3: fallback → 마크다운 저장
    return _render_markdown(spec, output_dir, filename)


def _find_korean_font() -> Optional[str]:
    """시스템에서 한글 폰트 찾기."""
    candidates = [
        "C:/Windows/Fonts/malgun.ttf",       # 맑은 고딕
        "C:/Windows/Fonts/NanumGothic.ttf",
        "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
        "/System/Library/Fonts/AppleSDGothicNeo.ttc",
    ]
    for p in candidates:
        if os.path.exists(p):
            return p
    return None


def _markdown_to_html(md: str, title: str = "Document") -> str:
    """마크다운 → 간단한 HTML (weasyprint용)."""
    # 간단한 마크다운 → HTML 변환
    import re
    html_body = md
    html_body = re.sub(r'^### (.+)$', r'<h3>\1</h3>', html_body, flags=re.MULTILINE)
    html_body = re.sub(r'^## (.+)$', r'<h2>\1</h2>', html_body, flags=re.MULTILINE)
    html_body = re.sub(r'^# (.+)$', r'<h1>\1</h1>', html_body, flags=re.MULTILINE)
    html_body = re.sub(r'^- (.+)$', r'<li>\1</li>', html_body, flags=re.MULTILINE)
    html_body = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', html_body)
    html_body = html_body.replace("\n\n", "</p><p>").replace("\n", "<br>")

    return f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>{title}</title>
<style>
body {{ font-family: 'Malgun Gothic', 'Nanum Gothic', sans-serif; max-width: 800px; margin: 40px auto; padding: 0 20px; line-height: 1.6; }}
h1 {{ border-bottom: 2px solid #333; padding-bottom: 8px; }}
h2 {{ color: #2c3e50; margin-top: 24px; }}
li {{ margin: 4px 0; }}
</style></head><body><p>{html_body}</p></body></html>"""


# ═══════════════════════════════════════════
#  Doc/README → .md
# ═══════════════════════════════════════════

def _render_markdown(spec: dict, output_dir: str, filename: str, ext: str = "md") -> str:
    """JSON spec → .md 파일."""
    md = _spec_to_markdown(spec)
    path = os.path.join(output_dir, f"{filename}.{ext}")
    with open(path, "w", encoding="utf-8") as f:
        f.write(md)
    return path


def _spec_to_markdown(spec: dict) -> str:
    """JSON spec → 마크다운 문자열."""
    lines = []
    title = spec.get("document_title", spec.get("deck_title", "Document"))
    lines.append(f"# {title}\n")

    # doc/pdf: sections
    for section in spec.get("sections", []):
        heading = section.get("heading", "")
        content = section.get("content", "")
        if heading:
            lines.append(f"\n## {heading}\n")
        if content:
            lines.append(content)
        for kp in section.get("key_points", []):
            lines.append(f"- {kp}")

    # ppt: slides (마크다운 형태)
    for slide in spec.get("slides", []):
        lines.append(f"\n## {slide.get('title', '')}\n")
        for b in slide.get("bullets", []):
            lines.append(f"- {b}")
        if slide.get("content"):
            lines.append(slide["content"])

    return "\n".join(lines)
